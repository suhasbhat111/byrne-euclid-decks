#!/usr/bin/env python3
"""
Build basic_principles.pptx from the 5 existing Basic Principles PDFs.
Each PDF page becomes one slide (full-page image, centred).
Canonical order: definitions, postulates, axioms, elucidations, symbols.

Usage: python3 basic_principles_pptx.py
Output: ../build/basic_principles.pptx
"""
import re, subprocess, tempfile, os
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

BUILD = Path(__file__).resolve().parent.parent / "build"
OUT   = BUILD / "basic_principles.pptx"
PDFTOCAIRO = "/opt/homebrew/bin/pdftocairo"

PDFS = [
    BUILD / "basic_principles_definitions.pdf",
    BUILD / "basic_principles_postulates.pdf",
    BUILD / "basic_principles_axioms.pdf",
    BUILD / "basic_principles_elucidations.pdf",
    BUILD / "basic_principles_symbols.pdf",
]

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

with tempfile.TemporaryDirectory() as tmp:
    for pdf in PDFS:
        info = subprocess.run(['pdfinfo', str(pdf)], capture_output=True, text=True).stdout
        npages = int(re.search(r'Pages:\s+(\d+)', info).group(1))
        for p in range(1, npages + 1):
            stem = os.path.join(tmp, f"{pdf.stem}_p{p}")
            subprocess.run([
                PDFTOCAIRO, '-png', '-r', '200',
                '-f', str(p), '-l', str(p), '-singlefile', str(pdf), stem
            ], capture_output=True)
            png = stem + '.png'
            if not os.path.exists(png):
                print(f"  WARN: could not render {pdf.stem} p{p}")
                continue
            img = PILImage.open(png)
            w_in, h_in = img.width / 200.0, img.height / 200.0
            sw, sh = SLIDE_W / Inches(1), SLIDE_H / Inches(1)
            scale = min(sw / w_in, sh / h_in)
            fw, fh = w_in * scale, h_in * scale
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(png,
                                     Inches((sw - fw) / 2), Inches((sh - fh) / 2),
                                     width=Inches(fw), height=Inches(fh))
            print(f"  {pdf.stem} p{p}")

prs.save(str(OUT))
print(f"Saved → {OUT}  ({len(prs.slides)} slides)")
