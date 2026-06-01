#!/usr/bin/env python3
"""
Merge all individual decks into build/book_I_complete.pptx.
Order: basic_principles, then I01–I48 (each 2–3 slides).

Usage: python3 merge_pptx.py
"""
import io, copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

BUILD = Path(__file__).resolve().parent.parent / "build"
R_EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'

SOURCES = [BUILD / "basic_principles.pptx"] + \
          [BUILD / f"I{n:02d}.pptx" for n in range(1, 49)]

dest = Presentation()
dest.slide_width  = Inches(13.33)
dest.slide_height = Inches(7.50)
blank = dest.slide_layouts[6]

total = 0
for src_path in SOURCES:
    if not src_path.exists():
        print(f"  skip (missing): {src_path.name}")
        continue
    src = Presentation(str(src_path))
    for src_slide in src.slides:
        dest_slide = dest.slides.add_slide(blank)
        rId_map = {}
        for rId, rel in src_slide.part.rels.items():
            if 'image' not in rel.reltype:
                continue
            blob = rel.target_part.blob
            before = set(dest_slide.part.rels.keys())
            pic = dest_slide.shapes.add_picture(io.BytesIO(blob), Emu(0), Emu(0), Emu(1), Emu(1))
            new = set(dest_slide.part.rels.keys()) - before
            if new:
                rId_map[rId] = new.pop()
            pic._element.getparent().remove(pic._element)
        dest_tree = dest_slide.shapes._spTree
        for child in list(dest_tree):
            dest_tree.remove(child)
        for child in src_slide.shapes._spTree:
            elem = copy.deepcopy(child)
            for el in elem.iter():
                old = el.get(R_EMBED)
                if old and old in rId_map:
                    el.set(R_EMBED, rId_map[old])
            dest_tree.append(elem)
        total += 1

out = BUILD / "book_I_complete.pptx"
dest.save(str(out))
print(f"Saved → {out}  ({total} slides)")
