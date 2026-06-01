#!/usr/bin/env python3
"""
packager.py  —  Byrne generator .tex  →  editable PPTX

Proof text is real python-pptx text runs (EB Garamond).
Byrne symbols (lines, angles, polygons, circles, inline figures) are
small PNGs placed at computed absolute (x, y) positions derived from
EB Garamond font metrics.

Usage:  python3 packager.py ../generators/generator_I01.tex
Output: ../build/I01.pptx
"""

import hashlib, os, re, sys, subprocess, tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Tuple, Optional, Dict

from PIL import Image as PILImage, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

# ── Paths ─────────────────────────────────────────────────────────────────────

LUALATEX  = "/usr/local/texlive/2026/bin/universal-darwin/lualatex"
GS        = "/opt/homebrew/bin/gs"
RENDER_DPI = 150   # px/in for symbol renders; 150 is enough for crisp inline symbols

FONT_DIR  = "/usr/local/texlive/2026/texmf-dist/fonts/opentype/public/ebgaramond/"
FONT_REG  = FONT_DIR + "EBGaramond-Regular.otf"
FONT_ITAL = FONT_DIR + "EBGaramond-Italic.otf"
FONT_BOLD = FONT_DIR + "EBGaramond-Bold.otf"

# ── Slide geometry (inches as floats; EMU wrappers below) ───────────────────────
# Design goal: fill the slide.  Generous margins, a confident title with a Byrne
# tricolor accent rule, a large vertically-centred figure, and a text column that
# grows its font + leading to occupy the full body height (fit-to-fill).

SW = 13.333        # slide width  (16:9)
SH = 7.5           # slide height
MARGIN_F = 0.55    # outer margin

TITLE_Y_F   = 0.38         # title baseline-ish top
TITLE_PT    = 27           # title font size (pt)
RULE_Y_F    = 1.06         # Byrne tricolour accent rule under the title
RULE_H_F    = 0.045        # rule thickness

CONTENT_TOP_F = 1.34       # body starts below the rule
CONTENT_BOT_F = 7.06       # body ends above bottom margin
CONTENT_H_F   = CONTENT_BOT_F - CONTENT_TOP_F

COL_GAP_F  = 0.55
FIG_LEFT_F = MARGIN_F
FIG_W_F    = 5.55                                   # figure column width
TEXT_LEFT_F = FIG_LEFT_F + FIG_W_F + COL_GAP_F
TEXT_W_F    = SW - TEXT_LEFT_F - MARGIN_F          # ≈ 6.13 in

# Font-size envelopes (pt): (base, max, min) — the engine picks the largest that
# fits, then expands leading + vertically centres to fill the body.
STMT_FONTS  = (18, 27, 13)
PROOF_FONTS = (15, 23, 11)

LEAD_MIN = 1.42    # tightest acceptable line-height multiplier
LEAD_MAX = 2.55    # most generous (used when content is short, to fill height)
FILL_TGT = 0.94    # fraction of body height the text block should occupy

# ── EMU wrappers (python-pptx needs EMU) ────────────────────────────────────────
SLIDE_W = Inches(SW)
SLIDE_H = Inches(SH)

# Byrne palette (approximated from the 1847 colour plates)
INK    = RGBColor(0x22, 0x22, 0x22)
BY_RED = RGBColor(0xC0, 0x3A, 0x22)
BY_BLU = RGBColor(0x27, 0x55, 0x97)
BY_YEL = RGBColor(0xE6, 0xA5, 0x17)

# Back-compat aliases (a few helpers still reference these)
MARGIN     = Inches(MARGIN_F)
CONTENT_Y  = Inches(CONTENT_TOP_F)
CONTENT_H  = Inches(CONTENT_H_F)
TEXT_LEFT  = Inches(TEXT_LEFT_F)
TEXT_MAX_W = Inches(TEXT_W_F)
TITLE_Y    = Inches(TITLE_Y_F)
TITLE_H    = Inches(0.7)
FIG_LEFT   = Inches(FIG_LEFT_F)
FIG_MAX_W  = Inches(FIG_W_F)

STMT_PT   = 16
PROOF_PT  = 14
LINE_MULT = 1.55   # legacy default (fit-to-fill computes its own leading)

# ── Token types ───────────────────────────────────────────────────────────────

@dataclass
class Tok:
    """Logical unit from the LaTeX scanner."""
    kind: str        # 'text' | 'pic' | 'space' | 'br'
    text: str = ""   # plain text (kind='text' or 'space')
    call: str = ""   # full LaTeX call to render (kind='pic')
    bold: bool = False
    italic: bool = False
    voffset_pt: float = 0.0   # baseline raise in points (for tall offsetPicture pics)
    large: bool = False       # True for offsetPicture / drawFromCurrentPicture
    br_skip_pt: float = 0.0  # extra skip after line break


# Macros that must be rendered as images
PIC_MACROS = {
    'drawUnitLine', 'drawProportionalLine', 'drawSizedLine',
    'drawLine', 'drawAngle', 'drawAngleWithSides', 'drawPolygon',
    'drawCircle', 'drawArc', 'drawPoint', 'drawPointM',
    'drawRightAngle', 'drawTwoRightAngles',
    'offsetPicture', 'drawFromCurrentPicture', 'defineFromCurrentPicture',
}

# ── LaTeX templates ───────────────────────────────────────────────────────────

def _pic_tex(scale: str, body: str, defs: str, call: str,
             large: bool = False) -> str:
    """Standalone doc that renders one pictogram in the full figure context."""
    # Large border for \offsetPicture / \drawFromCurrentPicture circles that
    # extend well above the text baseline; small border for unit lines/angles.
    border = "2cm" if large else "6pt"
    return (
        r"\documentclass[border=" + border + r"]{standalone}" + "\n"
        r"\usepackage[lining]{ebgaramond}" + "\n"
        r"\usepackage{amssymb}" + "\n"
        r"\usepackage{byrne}" + "\n"
        r"\def\mpPre{u := 1cm; textLabels := false;}" + "\n"
        r"\begin{document}" + "\n"
        r"\defineNewPicture[" + scale + r"]{%" + "\n"
        + body + "%\n"
        r"}" + "\n"
        + defs + "\n"
        + call + "\n"
        r"\end{document}" + "\n"
    )


def _figure_tex(scale: str, body: str) -> str:
    """Standalone doc that renders the full main figure."""
    return (
        r"\documentclass[border=1cm]{standalone}" + "\n"
        r"\usepackage{amssymb}" + "\n"
        r"\usepackage{byrne}" + "\n"
        r"\def\mpPre{u := 1cm; textLabels := false;}" + "\n"
        r"\begin{document}" + "\n"
        r"\defineNewPicture[" + scale + r"]{%" + "\n"
        + body + "%\n"
        r"}" + "\n"
        r"\drawCurrentPicture" + "\n"
        r"\end{document}" + "\n"
    )


# ── Renderer ──────────────────────────────────────────────────────────────────

def _trim_to_ink(path: str, pad_px: int = 3) -> None:
    """Crop a rendered PNG to its non-white ink bounding box (+ small padding).

    Used for `large` symbols (offsetPicture circles / inline figures): they are
    rendered with a generous 2cm border so the raised top is never cropped (the
    recurring 'chopped circle' bug), then trimmed here so the visible glyph fills
    the PNG and height-based scaling produces a correctly-sized symbol.  NOT used
    for thin line symbols, where trimming would destroy the aspect ratio.
    """
    from PIL import ImageChops
    img = PILImage.open(path).convert("RGB")
    white = PILImage.new("RGB", img.size, (255, 255, 255))
    bbox = ImageChops.difference(img, white).getbbox()
    if not bbox:
        return
    l, t, r, b = bbox
    l = max(0, l - pad_px); t = max(0, t - pad_px)
    r = min(img.width,  r + pad_px); b = min(img.height, b + pad_px)
    img.crop((l, t, r, b)).save(path)


def _render_png(tex: str, out: str, dpi: int = RENDER_DPI,
                trim: bool = False) -> bool:
    with tempfile.TemporaryDirectory() as d:
        t = os.path.join(d, "x.tex")
        p = os.path.join(d, "x.pdf")
        open(t, 'w').write(tex)
        r = subprocess.run(
            [LUALATEX, "-interaction=nonstopmode", "-output-directory", d, t],
            capture_output=True, text=True, cwd=d)
        if not os.path.exists(p):
            print("  lualatex failed:")
            for line in r.stdout.splitlines()[-15:]:
                print("   ", line)
            return False
        subprocess.run([
            GS, "-dNOPAUSE", "-dBATCH", "-sDEVICE=png16m",
            f"-r{dpi}", "-dFirstPage=1", "-dLastPage=1",
            f"-sOutputFile={out}", p
        ], capture_output=True)
        if os.path.exists(out) and trim:
            _trim_to_ink(out)
        return os.path.exists(out)


def cached_render(scale: str, body: str, defs: str, call: str,
                  cache: str, large: bool = False) -> Optional[str]:
    # Bump the cache key (v2) so stale pre-trim renders are not reused.
    h = hashlib.md5(("v2" + scale + body[:100] + defs + call + str(large)).encode()).hexdigest()
    out = os.path.join(cache, f"p_{h}.png")
    if not os.path.exists(out):
        # Trim large symbols (circles/figures) to ink so the 2cm safety border
        # doesn't shrink the visible glyph.
        if not _render_png(_pic_tex(scale, body, defs, call, large=large), out,
                           trim=large):
            return None
    return out


# ── Font metrics ──────────────────────────────────────────────────────────────

_fonts: Dict = {}

def _font(pt: float, bold: bool, italic: bool) -> ImageFont.FreeTypeFont:
    k = (pt, bold, italic)
    if k not in _fonts:
        path = FONT_BOLD if bold else (FONT_ITAL if italic else FONT_REG)
        _fonts[k] = ImageFont.truetype(path, int(pt * RENDER_DPI / 72))
    return _fonts[k]


def tw(text: str, pt: float, bold: bool = False, italic: bool = False) -> float:
    """Text width in inches at given point size."""
    if not text:
        return 0.0
    bb = _font(pt, bold, italic).getbbox(text)
    return (bb[2] - bb[0]) / RENDER_DPI


def cap_h_in(pt: float) -> float:
    """Cap height in inches (EB Garamond ≈ 0.72× point size)."""
    return pt * 0.72 / 72.0


def png_wh(path: str) -> Tuple[float, float]:
    img = PILImage.open(path)
    return img.width / RENDER_DPI, img.height / RENDER_DPI


# ── Parser ────────────────────────────────────────────────────────────────────

def _braced(s: str, pos: int) -> Tuple[str, int]:
    if pos >= len(s) or s[pos] != '{':
        return '', pos
    d, i = 1, pos + 1
    while i < len(s) and d:
        if s[i] == '{':   d += 1
        elif s[i] == '}': d -= 1
        i += 1
    return s[pos+1:i-1], i


def _optional(s: str, pos: int) -> Tuple[str, int]:
    if pos < len(s) and s[pos] == '[':
        e = s.index(']', pos)
        return s[pos+1:e], e + 1
    return '', pos


def _clean_title(s: str) -> str:
    """Strip LaTeX markup from title → plain text (B1)."""
    s = s.replace('~', ' ')
    s = re.sub(r'\\(?:enspace|quad|qquad|thinspace)\b', ' ', s)
    s = re.sub(r'\\[,;]\s*', ' ', s)
    s = re.sub(r'\\ ', ' ', s)
    s = re.sub(r'\\\.', '.', s)
    s = re.sub(r'\\[a-zA-Z]+\s*', '', s)
    return re.sub(r'\s+', ' ', s).strip()


def _clean_stmt(s: str) -> str:
    """Isolate prose from statement region — strips all structural LaTeX (B2, B5)."""
    # Cut at first \end{minipage}: removes figure-minipage leakage in PROB-tagged
    s = re.sub(r'\\end\{minipage\}.*', '', s, flags=re.DOTALL)
    # B5: font commands glued directly to next word (e.g. \itshapeFrom).
    # (?![a-z]) so we don't clip a longer command (e.g. never touch \centering's
    # stem inside another macro); uppercase/space/punct after the name still match.
    s = re.sub(r'\\(?:itshape|upshape|noindent|centering)(?![a-z])', ' ', s)
    # Strip structural commands
    s = re.sub(r'\\begin\{minipage\}[^\{]*\{[^}]*\}', '', s)
    s = re.sub(r'\\drawCurrentPicture\b[^\n]*\n?', '', s)
    s = re.sub(r'\\textbf\{[^}]+\}\s*', '', s)
    s = re.sub(r'\\vspace\{[^}]*\}\s*', '', s)
    s = re.sub(r'\\hfill\b\s*', '', s)
    s = re.sub(r'\\noindent\b\s*', '', s)
    s = re.sub(r'\\centering\b\s*', '', s)
    s = re.sub(r'\\medskip\b\s*', '', s)
    s = re.sub(r'\\bigskip\b\s*', '', s)
    s = re.sub(r'\\itshape\b\s*', '', s)
    s = re.sub(r'\\upshape\b\s*', '', s)
    s = re.sub(r'\\linewidth\b', '', s)
    # Final sweep: remove any remaining \begin{} / \end{}
    s = re.sub(r'\\begin\{[^}]+\}|\\end\{[^}]+\}', '', s)
    # Collapse spaces/tabs ONLY — newlines must survive because pic macros use
    # `%` line-continuation comments (e.g. \drawFromCurrentPicture{%\n ... \n}).
    # Collapsing the newline after `%` would comment out the closing brace and
    # break the render ("File ended while scanning ..."). The tokenizer treats
    # surviving newlines as ordinary spaces in prose, so wording is unaffected.
    s = re.sub(r'[ \t]+', ' ', s)
    s = re.sub(r'\n[ \t]*\n+', '\n', s)
    return s.strip()


def _clean_proof(s: str) -> str:
    """Strip structural LaTeX from proof region."""
    s = re.sub(r'\\begin\{center\}\s*', '', s)
    s = re.sub(r'\\end\{center\}\s*', '', s)
    s = re.sub(r'\\begin\{flushright\}.*?\\end\{flushright\}',
               '\n\nQ.~E.~D.', s, flags=re.DOTALL)
    s = re.sub(r'\\end\{document\}.*', '', s, flags=re.DOTALL)
    s = re.sub(r'\\end\{minipage\}', '', s)  # strip tag only; Q.E.D. may follow
    # B5: font/size commands glued to next word (e.g. \footnotesizeOn → " On").
    # (?![a-z]) prevents matching a prefix of a longer command — crucially
    # \small must NOT match inside \smallskip (which would leave a literal "skip").
    s = re.sub(r'\\(?:itshape|upshape|noindent|centering|footnotesize|small'
               r'|normalsize|large|Large|tiny|scriptsize)(?![a-z])', ' ', s)
    s = re.sub(r'\\(?:upshape|noindent|footnotesize|small|normalsize)\b\s*', '', s)
    # Vertical-skip macros that survived → drop (they carry no inline text)
    s = re.sub(r'\\(?:smallskip|medskip|bigskip)\b', '', s)
    return s.strip()


def _extract_picture(src: str, start: int = 0) -> Tuple[str, str, str, int]:
    """Extract (scale_a, scale_b, body, end_pos) from \\defineNewPicture at/after start.

    Handles both single-arg [sfA] and two-arg [sfA][sfB] forms (B4).
    scale_a = inline symbol scale; scale_b = main figure scale.
    """
    m = re.search(r'\\defineNewPicture\s*', src[start:])
    if not m:
        return '1/2', '1/2', '', start
    i = start + m.end()
    scales: List[str] = []
    while i < len(src) and src[i] == '[':
        j = src.index(']', i)
        scales.append(src[i+1:j])
        i = j + 1
    scale_a = scales[0] if scales else '1/2'
    scale_b = scales[1] if len(scales) > 1 else scale_a
    while i < len(src) and src[i] in ' \t\n':
        i += 1
    if i >= len(src) or src[i] != '{':
        return scale_a, scale_b, '', i
    d, i = 1, i + 1
    body_start = i
    while i < len(src) and d:
        if src[i] == '{':   d += 1
        elif src[i] == '}': d -= 1
        i += 1
    return scale_a, scale_b, src[body_start:i-1], i


def parse_generator(path: str) -> dict:
    src = open(path).read()
    prop_id = re.search(r'generator_(\w+)\.tex', path).group(1)

    # B4: extract figure body — supports one or two scale args
    scale_a, scale_b, body, _ = _extract_picture(src)

    # B1: title from first \textbf{...} in document
    tm = re.search(r'\\textbf\{([^}]+)\}', src)
    title = _clean_title(tm.group(1)) if tm else prop_id

    # All 48 generators have exactly one \itshape — statement starts there
    its_m   = re.search(r'\\itshape', src)
    its_pos = its_m.end() if its_m else len(src)

    # Proof-start marker priority: %PROOF > \upshape > \begin{center} after \itshape (B3)
    pm      = re.search(r'%\s*PROOF', src)
    ups_m   = re.search(r'\\upshape', src[its_pos:]) if its_m else None
    if ups_m and its_m:
        ups_m = re.compile(r'\\upshape').search(src, its_pos)
    ctr_m   = re.compile(r'\\begin\{center\}').search(src, its_pos + 1) if its_m else None
    minip_m = re.compile(r'\\end\{minipage\}').search(src, its_pos) if its_m else None

    if pm:
        # PROB-tagged I01–I10
        stmt_end    = pm.start()
        proof_start = pm.end()
    elif ups_m:
        # THEOR with \upshape (35 generators)
        stmt_end    = ups_m.start()
        proof_start = ups_m.start()
    else:
        # PROB-without-tag I11,I22,I23,I31,I42,I44,I45,I46
        # Statement ends at the first \end{minipage} (closes the statement minipage)
        # Proof starts at the \begin{center} that follows
        stmt_end    = minip_m.start() if minip_m else (ctr_m.start() if ctr_m else len(src))
        proof_start = ctr_m.start()   if ctr_m   else len(src)

    statement = _clean_stmt(src[its_pos:stmt_end])

    # I26 special case: \newpage splits two proof cases
    newpage_pos = src.find('\\newpage', proof_start)
    has_newpage = newpage_pos >= 0

    _re_minipend  = re.compile(r'\\end\{minipage\}')
    _re_upshape   = re.compile(r'\\upshape')
    # Terminator: include whole flushright block (Q.E.D.) via named group
    _re_proof_end = re.compile(
        r'(?P<fr>\\begin\{flushright\}.*?\\end\{flushright\})'
        r'|\\end\{minipage\}|\\end\{document\}',
        re.DOTALL)
    _re_fr = re.compile(r'\\begin\{flushright\}.*?\\end\{flushright\}', re.DOTALL)

    def _slice_proof(start: int) -> str:
        """Extract proof_raw from start, including Q.E.D. flushright block if present."""
        m = _re_proof_end.search(src, start)
        if not m:
            return src[start:]
        if m.group('fr'):
            return src[start:m.end()]   # include the whole flushright block
        # minipage or document terminator — also check for flushright within next 80 chars
        fr_m = _re_fr.search(src, m.end(), m.end() + 80)
        if fr_m:
            return src[start:fr_m.end()]
        return src[start:m.start()]

    if has_newpage:
        # Case I: proof_start → first \end{minipage} before \newpage
        ci_end_m  = _re_minipend.search(src, proof_start)
        ci_end    = ci_end_m.start() if ci_end_m else newpage_pos
        proof_raw = src[proof_start:ci_end]

        # Case II: second \upshape after \newpage; Q.E.D. follows the closing \end{minipage}
        ups2_m = _re_upshape.search(src, newpage_pos)
        if ups2_m:
            proof2_raw = _slice_proof(ups2_m.start())
        else:
            proof2_raw = ''

        _, _, body2, _ = _extract_picture(src, start=newpage_pos)
    else:
        proof_raw  = _slice_proof(proof_start)
        proof2_raw = ''
        body2      = ''

    proof  = _clean_proof(proof_raw)
    proof2 = _clean_proof(proof2_raw) if proof2_raw else ''
    proof_centered = bool(re.search(r'\\begin\{center\}', proof_raw))

    return dict(
        prop_id=prop_id, title=title,
        scale=scale_a, scale_b=scale_b,
        body=body, body2=body2,
        statement=statement,
        proof=proof, proof2=proof2,
        proof_centered=proof_centered,
    )


# ── Tokenizer ─────────────────────────────────────────────────────────────────

def tokenize(src: str) -> List[Tok]:
    """
    Scan LaTeX text into a flat list of Tok objects.
    Each word is a separate 'text' tok; spaces are explicit 'space' toks.
    """
    toks: List[Tok] = []
    acc  = ""
    bold = False
    ital = False
    i    = 0

    def flush(b=None, it=None):
        nonlocal acc
        if not acc:
            return
        toks.append(Tok('text', text=acc, bold=b if b is not None else bold,
                        italic=it if it is not None else ital))
        acc = ""

    def push_space():
        if toks and toks[-1].kind != 'space':
            toks.append(Tok('space', text=' '))

    while i < len(src):
        c = src[i]

        # ── Math mode $...$ ──────────────────────────────────────────────────
        if c == '$':
            flush()
            j = src.find('$', i + 1)
            if j < 0:
                j = len(src) - 1
            math = src[i+1:j]
            i = j + 1
            if re.search(r'\\(?:' + '|'.join(PIC_MACROS) + r'|therefore)\b', math):
                toks.append(Tok('pic', call='$' + math + '$'))
            else:
                plain = (math.replace(r'\therefore', '∴')
                             .replace(r'\,', '')
                             .replace(r'\quad', '  '))
                plain = re.sub(r'\\[a-zA-Z]+\s*', '', plain)
                plain = plain.replace('{', '').replace('}', '').strip()
                if plain:
                    toks.append(Tok('text', text=plain))
            continue

        # ── LaTeX comment: % to end of line + newline + next-line indent ─────
        # `describe%\n\offsetPicture` joins with NO space (line continuation).
        # Pic-macro internal `%` is already consumed by _braced, so any `%`
        # reaching here is a real prose comment.
        if c == '%':
            nl = src.find('\n', i)
            if nl < 0:
                i = len(src)
            else:
                i = nl + 1
                while i < len(src) and src[i] in ' \t':
                    i += 1
            continue

        # ── Space / newline ──────────────────────────────────────────────────
        if c in ' \t\n':
            flush()
            push_space()
            i += 1
            continue

        # ── Braces (structural, not content) ────────────────────────────────
        if c in '{}':
            i += 1
            continue

        # ── Tilde = non-breaking space ───────────────────────────────────────
        if c == '~':
            flush()
            push_space()
            i += 1
            continue

        # ── Non-backslash: accumulate ────────────────────────────────────────
        if c != '\\':
            acc += c
            i += 1
            continue

        # ── Backslash ────────────────────────────────────────────────────────

        # Line break: \\ or \\[Npt]
        if i + 1 < len(src) and src[i+1] == '\\':
            flush()
            i += 2
            skip = 0.0
            if i < len(src) and src[i] == '[':
                e = src.index(']', i)
                try:
                    skip = float(src[i+1:e].replace('pt', '').strip())
                except ValueError:
                    skip = 0.0
                i = e + 1
            toks.append(Tok('br', br_skip_pt=skip))
            continue

        # Collect command name
        j = i + 1
        while j < len(src) and src[j].isalpha():
            j += 1
        cmd = src[i+1:j]

        # ── Pictogram macros ─────────────────────────────────────────────────
        if cmd in PIC_MACROS:
            flush()
            ci   = j
            full = '\\' + cmd
            opts = []
            args = []
            while ci < len(src) and src[ci] == '[':
                opt, ci = _optional(src, ci)
                opts.append(opt)
                full += f'[{opt}]'
            while ci < len(src) and src[ci] == '{':
                arg, ci = _braced(src, ci)
                args.append(arg)
                full += '{' + arg + '}'
            voff  = 0.0
            large = False
            if cmd == 'offsetPicture':
                # \offsetPicture{<vraise>}{<hshift>}{<content>} — offset is the
                # FIRST BRACE arg (e.g. {15pt}), not a bracket option.  These are
                # raised circles/figures that need the large 2cm render border so
                # the raised top is not cropped (the recurring "chopped top" bug).
                large = True
                raw = (args[0] if args else (opts[0] if opts else '0'))
                try:
                    voff = float(raw.replace('pt', '').strip())
                except ValueError:
                    voff = 0.0
            elif cmd in ('drawFromCurrentPicture', 'defineFromCurrentPicture'):
                large = True
            toks.append(Tok('pic', call=full, voffset_pt=voff, large=large))
            i = ci
            continue

        # ── Named sub-figure back-references: \triangleABC, \pointA, etc. ───
        if re.match(r"(?:triangle|polygon|circle|arc|angles?|point|line|"
                    r"rect|square)[A-Z][A-Za-z0-9']*$", cmd):
            flush()
            toks.append(Tok('pic', call='\\' + cmd))
            i = j
            continue

        # ── Text-decoration commands ─────────────────────────────────────────
        if cmd in ('textbf', 'textit', 'emph', 'textsc', 'textsl'):
            flush()
            arg, ni = _braced(src, j)
            is_bold = (cmd == 'textbf')
            is_ital = cmd in ('textit', 'emph', 'textsl')
            sub = tokenize(arg)
            for t in sub:
                if is_bold: t.bold = True
                if is_ital: t.italic = True
            toks.extend(sub)
            i = ni
            continue

        # ── Inline math accents/macros that become text ──────────────────────
        if cmd == 'therefore':
            flush()
            toks.append(Tok('text', text='∴'))
            i = j
            continue

        if cmd == 'byref':
            flush()
            arg, ni = _braced(src, j)
            toks.append(Tok('text', text=f'({arg})'))
            i = ni
            continue

        if cmd == 'quad':
            flush()
            toks.append(Tok('space', text='   '))
            i = j
            continue

        if cmd in ('thinspace', ','):
            flush()
            push_space()
            i = j
            continue

        if cmd in ('medskip',):
            flush()
            toks.append(Tok('br', br_skip_pt=6.0))
            i = j
            continue

        if cmd in ('bigskip',):
            flush()
            toks.append(Tok('br', br_skip_pt=12.0))
            i = j
            continue

        if cmd in ('noindent', 'centering', 'small', 'normalsize',
                   'large', 'Large', 'footnotesize', 'par'):
            i = j
            continue

        # ── \begin{...} and \end{...}: consume the {arg} and discard ─────────
        if cmd in ('begin', 'end'):
            flush()
            _, ni = _braced(src, j)   # consume e.g. {document} {center} etc.
            i = ni
            continue

        # ── Escaped chars: \. \ \, \; etc. ──────────────────────────────────
        if j == i + 1:
            ch = src[i+1] if i+1 < len(src) else ''
            if ch == ' ':
                flush()
                push_space()
                i += 2
            elif ch in '.,;:!?-~':
                acc += ch
                i += 2
            else:
                i += 2
            continue

        # Unknown command: skip
        i = j

    flush()
    return toks


# ── Layout engine ─────────────────────────────────────────────────────────────

@dataclass
class LTok:
    """Layout token: Tok + computed pixel-independent dimensions."""
    kind: str
    text:   str   = ""
    bold:   bool  = False
    italic: bool  = False
    w:      float = 0.0    # width in inches
    h:      float = 0.0    # height in inches
    png:    str   = ""     # PNG path (pic tokens)
    voff:   float = 0.0    # vertical offset in inches (large pics)
    large:  bool  = False
    br_skip: float = 0.0   # extra vertical skip after a 'br'


def collect_named_defs(*tok_lists) -> Dict[str, str]:
    """Scan pic tokens across regions for `[..][name]{..}` definitions.

    Returns name → defining call.  Built globally (statement + proof) so a
    back-reference like \\pointA defined in the statement still resolves when
    used in the proof.
    """
    named_defs: Dict[str, str] = {}
    for toks in tok_lists:
        for tok in toks:
            if tok.kind != 'pic':
                continue
            opts = re.findall(r'\[([^\]]*)\]', tok.call)
            if len(opts) >= 2 and opts[1].strip():
                named_defs[opts[1].strip()] = tok.call
    return named_defs


def build_layout(toks: List[Tok], scale: str, body: str,
                 font_pt: float, cache: str,
                 named_defs: Optional[Dict[str, str]] = None) -> List[LTok]:
    """Render every pic token; measure every text token.  Returns LToks."""

    # Pre-scan: find all named sub-figure defining calls (local fallback if no
    # global map was supplied).
    if named_defs is None:
        named_defs = collect_named_defs(toks)

    result: List[LTok] = []
    ch = cap_h_in(font_pt)

    for tok in toks:
        if tok.kind == 'text':
            result.append(LTok('text', text=tok.text,
                               bold=tok.bold, italic=tok.italic,
                               w=tw(tok.text, font_pt, tok.bold, tok.italic)))

        elif tok.kind == 'space':
            result.append(LTok('space', text=' ',
                               w=tw(' ', font_pt)))

        elif tok.kind == 'br':
            result.append(LTok('br', br_skip=tok.br_skip_pt / 72.0))

        elif tok.kind == 'pic':
            # Build preamble for named sub-figure back-references
            defs_latex = ""
            ref = re.match(r"\\((?:triangle|polygon|circle|arc|angles?|point|"
                           r"line|rect|square)[A-Z][A-Za-z0-9']*)$",
                           tok.call)
            if ref:
                name = ref.group(1)
                if name in named_defs:
                    defs_latex = f"\\setbox0=\\hbox{{{named_defs[name]}}}%\n"

            png = cached_render(scale, body, defs_latex, tok.call,
                                cache, large=tok.large)
            if png:
                w_in, h_in = png_wh(png)
                # Scale heights:
                #   large (circles/inline figures) → 3.0 × cap_h
                #   math group ($...$)             → 2.0 × cap_h (contains text + symbols)
                #   regular inline symbol          → 1.5 × cap_h
                is_math = tok.call.startswith('$')
                if tok.large:
                    target_h = ch * 2.5
                elif is_math:
                    target_h = ch * 2.0
                else:
                    target_h = ch * 1.5
                s = target_h / h_in if h_in > 0 else 1.0
                fw, fh = w_in * s, h_in * s
                # Vertical offset in inches (from tok baseline raise in pt)
                voff = tok.voffset_pt / 72.0 if tok.large else 0.0
                result.append(LTok('pic', w=fw, h=fh, png=png,
                                   voff=voff, large=tok.large))
            else:
                result.append(LTok('text', text='[?]',
                                   w=tw('[?]', font_pt)))

    return result


def flow_lines(ltoks: List[LTok], max_w: float) -> List[List[LTok]]:
    """Word-wrap a list of LToks into lines that fit within max_w inches."""
    lines: List[List[LTok]] = []
    cur:   List[LTok] = []
    cur_w = 0.0

    def commit():
        nonlocal cur, cur_w
        # strip trailing spaces
        while cur and cur[-1].kind == 'space':
            cur.pop()
        lines.append(cur)
        cur = []
        cur_w = 0.0

    for lt in ltoks:
        if lt.kind == 'br':
            commit()
            lines.append([lt])   # spacer sentinel
            continue

        # Skip leading spaces on a new line
        if lt.kind == 'space' and not cur:
            continue

        # Overflow check (not for spaces — they wrap with the next word)
        if lt.kind != 'space' and cur and cur_w + lt.w > max_w:
            commit()

        cur.append(lt)
        cur_w += lt.w

    if cur:
        commit()

    return lines


# ── PPTX placement ────────────────────────────────────────────────────────────

def _add_text_run(slide, text: str, x: float, y: float,
                  line_h: float, pt: float, bold: bool, italic: bool):
    """Place one word as a zero-margin text box."""
    w = tw(text, pt, bold, italic) + Inches(0.04) / Inches(1)
    txb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                   Inches(w), Inches(line_h))
    tf = txb.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text          = text
    run.font.name     = "EB Garamond"
    run.font.size     = Pt(pt)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def _flush_phrase(slide, buf: List[LTok], bx: float, y: float,
                  line_h: float, font_pt: float):
    """Emit one multi-run textbox for a run of consecutive text/space LToks.

    Uses SHAPE_TO_FIT_TEXT so the box auto-expands to fit whatever font the
    viewer uses — prevents text clipping when EB Garamond is unavailable and a
    wider font is substituted.  The box background is transparent (no fill) so
    auto-expansion never obscures adjacent symbol images.
    """
    if not buf:
        return
    total = sum(lt.w for lt in buf)
    if total <= 0:
        return
    # Initial width = EB Garamond measurement (the box will auto-expand if needed)
    w_box = total * 1.25 + Inches(0.15) / Inches(1)
    txb = slide.shapes.add_textbox(Inches(bx), Inches(y),
                                   Inches(w_box), Inches(line_h))
    tf = txb.text_frame
    tf.word_wrap = False
    # Auto-expand box width to fit rendered text regardless of font substitution.
    # Textbox background is transparent so expansion cannot cover symbol images.
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Emu(0)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    for lt in buf:
        run = para.add_run()
        run.text           = lt.text
        run.font.name      = "EB Garamond"
        run.font.size      = Pt(font_pt)
        run.font.bold      = lt.bold
        run.font.italic    = lt.italic
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def place_line(slide, ltoks: List[LTok], x0: float, y: float,
               line_h: float, font_pt: float,
               centered: bool = False, col_w: float = 0.0):
    """Place one line of LToks — consecutive text/space runs coalesced into one textbox."""
    if not ltoks or (len(ltoks) == 1 and ltoks[0].kind == 'br'):
        return

    total_w = sum(lt.w for lt in ltoks)
    x = x0
    if centered and col_w > 0 and total_w < col_w:
        x = x0 + (col_w - total_w) / 2.0

    baseline = y + line_h * 0.78
    ch = cap_h_in(font_pt)

    phrase_buf: List[LTok] = []
    phrase_x = x

    for lt in ltoks:
        if lt.kind == 'br':
            break
        if lt.kind in ('text', 'space'):
            phrase_buf.append(lt)
            x += lt.w
        elif lt.kind == 'pic':
            _flush_phrase(slide, phrase_buf, phrase_x, y, line_h, font_pt)
            phrase_buf = []
            if lt.large:
                img_y = baseline - lt.h + lt.voff
            else:
                img_y = y + (line_h - lt.h) * 0.35
            slide.shapes.add_picture(
                lt.png,
                Inches(x), Inches(img_y),
                width=Inches(lt.w), height=Inches(lt.h)
            )
            x += lt.w
            phrase_x = x

    _flush_phrase(slide, phrase_buf, phrase_x, y, line_h, font_pt)


def _add_rule(slide, x, y, w, h, color):
    """Thin filled rectangle (no outline) used for the title accent rule."""
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_title(slide, prop_title: str, label: str):
    """Confident title + a Byrne tricolour accent rule beneath it."""
    txb = slide.shapes.add_textbox(Inches(MARGIN_F), Inches(TITLE_Y_F),
                                   Inches(SW - 2 * MARGIN_F), Inches(0.7))
    tf  = txb.text_frame
    tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = Emu(0)
    p   = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = prop_title
    r1.font.size = Pt(TITLE_PT); r1.font.bold = True
    r1.font.name = "EB Garamond"; r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = f"   {label}"
    r2.font.size = Pt(TITLE_PT - 5); r2.font.bold = False; r2.font.italic = True
    r2.font.name = "EB Garamond"; r2.font.color.rgb = BY_RED

    # Tricolour rule (red / blue / yellow) — a nod to Byrne's colour method.
    full = SW - 2 * MARGIN_F
    seg  = full / 3.0
    for i, col in enumerate((BY_RED, BY_BLU, BY_YEL)):
        _add_rule(slide, MARGIN_F + i * seg, RULE_Y_F, seg, RULE_H_F, col)


def place_main_figure(slide, png_path: str):
    """Large figure, centred in the figure column AND vertically in the body."""
    if not os.path.exists(png_path):
        return
    w_in, h_in = png_wh(png_path)
    s = min(FIG_W_F / w_in, CONTENT_H_F / h_in)
    fw, fh = w_in * s, h_in * s
    x = FIG_LEFT_F + (FIG_W_F - fw) / 2.0
    y = CONTENT_TOP_F + (CONTENT_H_F - fh) / 2.0
    slide.shapes.add_picture(png_path, Inches(x), Inches(y),
                             width=Inches(fw), height=Inches(fh))


# ── Main assembler ────────────────────────────────────────────────────────────

def _block_height(lines: List[List[LTok]], font_pt: float, lead_mult: float,
                  break_scale: float = 1.0) -> float:
    """Total height (in) of flowed lines at a given leading multiplier."""
    lh = font_pt * lead_mult / 72.0
    total = 0.0
    for line in lines:
        if line and line[0].kind == 'br':
            base = line[0].br_skip if line[0].br_skip else lh * 0.35
            total += base * break_scale
        else:
            picmax = max((lt.h for lt in line if lt.kind == 'pic'), default=0.0)
            total += max(lh, picmax + font_pt * 0.20 / 72.0)
    return total


def fit_text_block(raw_text: str, scale: str, body: str, cache: str,
                   named_defs: Optional[Dict[str, str]],
                   col_w: float, avail_h: float, fonts: Tuple[int, int, int]):
    """Choose font size + leading so the text fills the body, then vertically
    centre it.  Returns (lines, per_line_heights, font_pt, top_y)."""
    base_pt, max_pt, min_pt = fonts
    toks = tokenize(raw_text)

    # 1. Largest font (≤max) whose content fits avail_h at the tight leading.
    chosen_pt, chosen_lines = min_pt, None
    for pt in range(max_pt, min_pt - 1, -1):
        ltoks = build_layout(toks, scale, body, pt, cache, named_defs=named_defs)
        lines = flow_lines(ltoks, col_w)
        if _block_height(lines, pt, LEAD_MIN) <= avail_h:
            chosen_pt, chosen_lines = pt, lines
            break
    if chosen_lines is None:                       # nothing fit → smallest font
        ltoks = build_layout(toks, scale, body, min_pt, cache, named_defs=named_defs)
        chosen_lines = flow_lines(ltoks, col_w)
        chosen_pt = min_pt

    # 2. Expand leading (and paragraph breaks) to fill FILL_TGT of the body.
    nat = _block_height(chosen_lines, chosen_pt, LEAD_MIN)
    target = FILL_TGT * avail_h
    n_text = sum(1 for ln in chosen_lines if not (ln and ln[0].kind == 'br'))
    lead = LEAD_MIN
    bscale = 1.0
    if nat < target and n_text > 0:
        # Solve leading so text-line growth covers most of the gap; let breaks
        # stretch a bit too for rhythm.
        extra = target - nat
        per_line = extra / n_text
        lead = min(LEAD_MAX, LEAD_MIN + per_line * 72.0 / chosen_pt)
        # If leading alone capped out, open up the paragraph breaks to help fill.
        rem = target - _block_height(chosen_lines, chosen_pt, lead)
        if rem > 0:
            bscale = min(2.6, 1.0 + rem / max(0.4, avail_h))

    # 3. Per-line heights + vertical centring.
    line_heights = []
    lh = chosen_pt * lead / 72.0
    for line in chosen_lines:
        if line and line[0].kind == 'br':
            base = line[0].br_skip if line[0].br_skip else lh * 0.35
            line_heights.append(base * bscale)
        else:
            picmax = max((lt.h for lt in line if lt.kind == 'pic'), default=0.0)
            line_heights.append(max(lh, picmax + chosen_pt * 0.20 / 72.0))
    block = sum(line_heights)
    top = CONTENT_TOP_F + max(0.0, (avail_h - block) / 2.0)
    return chosen_lines, line_heights, chosen_pt, top


def build_slide(prs, blank_layout, data: dict, label: str, raw_text: str,
                fonts: Tuple[int, int, int], scale_a: str, body: str,
                cache: str, named_defs, fig_png: str, centered: bool):
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, data['title'], label)
    place_main_figure(slide, fig_png)

    lines, line_heights, font_pt, top = fit_text_block(
        raw_text, scale_a, body, cache, named_defs,
        TEXT_W_F, CONTENT_H_F, fonts)

    y = top
    bottom = CONTENT_BOT_F + 0.05
    for line, lh in zip(lines, line_heights):
        if y > bottom:
            print(f"  OVERFLOW WARNING: {data['prop_id']} '{label}' — content clipped")
            break
        if line and line[0].kind == 'br':
            y += lh
            continue
        place_line(slide, line, TEXT_LEFT_F, y, lh, font_pt,
                   centered=centered, col_w=TEXT_W_F)
        y += lh


def build_pptx(data: dict, out_path: str, cache_dir: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    scale_a    = data['scale']                  # inline symbol scale
    scale_b    = data.get('scale_b', scale_a)   # main figure scale (B4)
    body       = data['body']
    body2      = data.get('body2', '') or body

    # Main figure — use scale_b for the standalone render (B4)
    fig_png = os.path.join(cache_dir, f"fig_{data['prop_id']}.png")
    if not os.path.exists(fig_png):
        print("  Rendering main figure...")
        _render_png(_figure_tex(scale_b, body), fig_png, dpi=200)

    # Global named-def map: a \pointA defined in the statement must resolve when
    # referenced in the proof, so collect across all regions first.
    g_defs = collect_named_defs(
        tokenize(data['statement']),
        tokenize(data['proof']),
        tokenize(data.get('proof2', '') or ''),
    )

    print("  Laying out statement...")
    build_slide(prs, blank, data, "Statement", data['statement'],
                STMT_FONTS, scale_a, body, cache_dir, g_defs, fig_png,
                centered=False)

    print("  Laying out proof...")
    proof_label = "Proof — Case 1" if data.get('proof2') else "Proof"
    build_slide(prs, blank, data, proof_label, data['proof'],
                PROOF_FONTS, scale_a, body, cache_dir, g_defs, fig_png,
                centered=data['proof_centered'])

    # I26: second case slide with its own figure
    if data.get('proof2'):
        fig2_png = os.path.join(cache_dir, f"fig_{data['prop_id']}_case2.png")
        if not os.path.exists(fig2_png):
            print("  Rendering Case II figure...")
            _render_png(_figure_tex(scale_b, body2), fig2_png, dpi=200)
        build_slide(prs, blank, data, "Proof — Case 2", data['proof2'],
                    PROOF_FONTS, scale_a, body2, cache_dir, g_defs, fig2_png,
                    centered=data['proof_centered'])

    prs.save(out_path)
    print(f"  Saved → {out_path}")


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2 or sys.argv[1] in ('-h', '--help'):
        print("Usage: packager.py <path/to/generator_I<NN>.tex>  |  packager.py --all")
        sys.exit(1)

    if sys.argv[1] == '--all':
        gen_dir   = Path(__file__).resolve().parent.parent / "generators"
        build_dir = Path(__file__).resolve().parent.parent / "build"
        build_dir.mkdir(exist_ok=True)
        cache_dir = str(build_dir / "pic_cache")
        os.makedirs(cache_dir, exist_ok=True)
        for tex_path in sorted(gen_dir.glob("generator_I*.tex")):
            data = parse_generator(str(tex_path))
            out_pptx = str(build_dir / f"{data['prop_id']}.pptx")
            print(f"Packaging {data['prop_id']}  ({data['title']})...")
            build_pptx(data, out_pptx, cache_dir)
        print("Done — all propositions packaged.")
        return

    tex_path = sys.argv[1]
    data     = parse_generator(tex_path)

    build_dir = Path(tex_path).resolve().parent.parent / "build"
    build_dir.mkdir(exist_ok=True)
    cache_dir = str(build_dir / "pic_cache")
    os.makedirs(cache_dir, exist_ok=True)

    out_pptx = str(build_dir / f"{data['prop_id']}.pptx")
    print(f"Packaging {data['prop_id']}  ({data['title']})...")
    build_pptx(data, out_pptx, cache_dir)
    print("Done.")


if __name__ == "__main__":
    main()
